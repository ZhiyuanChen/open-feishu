# OpenFeishu
# Copyright (C) 2024-Present  DanLing

# This file is part of OpenFeishu.

# OpenFeishu is free software: you can redistribute it and/or modify
# it under the terms of the GNU Affero General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# any later version.

# OpenFeishu is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU Affero General Public License for more details.

# You should have received a copy of the GNU Affero General Public License
# along with this program.  If not, see <http://www.gnu.org/licenses/>.

# For additional terms and clarifications, please refer to our License FAQ at:
# <https://multimolecule.danling.org/about/license-faq>.

r"""
附件内容提取：把图片 / PDF / Office / 文本附件的字节安全地抽取为「可交给多模态模型」的中性内容。

[feishu.attachments.extractor.AttachmentExtractor][] 抽象「字节 -> [feishu.attachments.extractor.ExtractedContent][]」，
默认实现 [feishu.attachments.extractor.SandboxedAttachmentExtractor][] 在一个可被 SIGKILL 终止的子进程中运行
提取，并施加体积、像素、ZIP（防 zip bomb）、页数、字符等上限，从而即便面对恶意构造的文件也不会拖垮或撑爆进程。
抽取结果是中性的（文本 + 图片 + 元信息，措辞为英文），不含任何产品提示词；产品侧可以使用
[feishu.attachments.analysis.analyze_attachment][] 的默认提示词，也可以直接用
[feishu.attachments.extractor.to_openai_content_parts][] 配上自己的分析提示词构造模型输入。

PDF / Office 解析依赖可选库（PyMuPDF、python-docx、openpyxl、python-pptx、Pillow），按需惰性导入；缺失时对应
格式优雅降级为「无法提取」而非报错。
"""

from __future__ import annotations

import asyncio
import base64
import mimetypes
import multiprocessing
import queue
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any, Mapping, Protocol, runtime_checkable

# Default guard limits (mirror battle-tested values). Override via ExtractLimits.
IMAGE_MAX_PIXELS = 16_000_000
INLINE_TEXT_MAX_CHARS = 40_000
EXTRACT_TEXT_MAX_CHARS = 80_000
PDF_EXTRACT_MAX_PAGES = 40
PDF_RENDER_MAX_PAGES = 4
PDF_RENDER_ZOOM = 1.7
PDF_RENDER_MAX_PIXELS = 16_000_000
PDF_RENDER_MAX_IMAGE_BYTES = 8 * 1024 * 1024
OFFICE_ZIP_MAX_ENTRIES = 300
OFFICE_ZIP_MAX_UNCOMPRESSED_BYTES = 40 * 1024 * 1024
DOCX_MAX_TABLES = 20
DOCX_MAX_TABLE_ROWS = 200
XLSX_MAX_SHEETS = 10
XLSX_MAX_ROWS_PER_SHEET = 200
PPTX_MAX_SLIDES = 80

_IMAGE_MIME_TYPES = {"image/jpeg", "image/png", "image/webp", "image/gif"}
_TEXT_MIME_TYPES = {
    "application/csv",
    "application/json",
    "application/xml",
    "application/x-yaml",
    "text/csv",
    "text/html",
    "text/markdown",
    "text/plain",
    "text/xml",
    "text/yaml",
}


@dataclass(frozen=True)
class ExtractLimits:
    r"""附件提取的各项上限；传给 [feishu.attachments.extractor.SandboxedAttachmentExtractor][] 以覆盖默认值。"""

    image_max_pixels: int = IMAGE_MAX_PIXELS
    image_max_bytes: int = 5 * 1024 * 1024
    inline_text_max_chars: int = INLINE_TEXT_MAX_CHARS
    extract_text_max_chars: int = EXTRACT_TEXT_MAX_CHARS
    pdf_extract_max_pages: int = PDF_EXTRACT_MAX_PAGES
    pdf_render_max_pages: int = PDF_RENDER_MAX_PAGES
    pdf_render_zoom: float = PDF_RENDER_ZOOM
    pdf_render_max_pixels: int = PDF_RENDER_MAX_PIXELS
    pdf_render_max_image_bytes: int = PDF_RENDER_MAX_IMAGE_BYTES
    office_zip_max_entries: int = OFFICE_ZIP_MAX_ENTRIES
    office_zip_max_uncompressed_bytes: int = OFFICE_ZIP_MAX_UNCOMPRESSED_BYTES
    docx_max_tables: int = DOCX_MAX_TABLES
    docx_max_table_rows: int = DOCX_MAX_TABLE_ROWS
    xlsx_max_sheets: int = XLSX_MAX_SHEETS
    xlsx_max_rows_per_sheet: int = XLSX_MAX_ROWS_PER_SHEET
    pptx_max_slides: int = PPTX_MAX_SLIDES


@dataclass
class ExtractedImage:
    r"""一张可交给多模态模型的图片：原始字节、媒体类型与（若可得）像素尺寸。"""

    data: bytes
    media_type: str
    width: int | None = None
    height: int | None = None


@dataclass
class ExtractedContent:
    r"""
    附件提取结果：分类、媒体类型、抽取文本、图片与元信息。

    `note` 为中性英文说明（如体积/像素超限、提取超时），不含任何产品提示词。

    Examples:
        >>> ExtractedContent(kind="text", media_type="text/plain", text="hi").kind
        'text'
    """

    kind: str
    media_type: str
    text: str | None = None
    images: list[ExtractedImage] = field(default_factory=list)
    total_pages: int | None = None
    truncated: bool = False
    note: str | None = None
    size_bytes: int = 0


@runtime_checkable
class AttachmentExtractor(Protocol):
    r"""
    附件提取协议：把附件字节抽取为 [feishu.attachments.extractor.ExtractedContent][]。

    内置实现为 [feishu.attachments.extractor.SandboxedAttachmentExtractor][]。该协议标注了 `runtime_checkable`。
    """

    async def extract(
        self, data: bytes, *, file_metadata: Mapping[str, Any], media_type: str | None = None
    ) -> ExtractedContent:
        r"""抽取附件内容；`media_type` 省略时按魔数 / 文件名推断。"""
        ...


class SandboxedAttachmentExtractor:
    r"""
    在可被 SIGKILL 终止的子进程中运行附件提取的默认 [feishu.attachments.extractor.AttachmentExtractor][]。

    超过 `max_bytes` 直接拒绝；其余在子进程中按 [feishu.attachments.extractor.ExtractLimits][] 提取，并以
    `timeout_seconds` 硬超时（子进程会被 terminate/kill），以信号量限制并发，从而抵御 zip bomb / 像素炸弹 /
    解析挂死等拒绝服务向量。任何异常都降级为带 `note` 的 `unknown` 结果，绝不抛出。

    Args:
        timeout_seconds: 单个附件的提取硬超时秒数。默认为 `20`。
        max_bytes: 允许提取的最大字节数。默认为 `16 MiB`。
        max_concurrency: 同时进行的提取数上限。默认为 `2`。
        limits: 各项提取上限 [feishu.attachments.extractor.ExtractLimits][]。

    Examples:
        >>> isinstance(SandboxedAttachmentExtractor(), AttachmentExtractor)
        True
    """

    def __init__(
        self,
        *,
        timeout_seconds: float = 20.0,
        max_bytes: int = 16 * 1024 * 1024,
        max_concurrency: int = 2,
        limits: ExtractLimits | None = None,
    ) -> None:
        self.timeout_seconds = timeout_seconds
        self.max_bytes = max_bytes
        self.limits = limits or ExtractLimits()
        self._semaphore = asyncio.Semaphore(max_concurrency)

    async def extract(
        self, data: bytes, *, file_metadata: Mapping[str, Any], media_type: str | None = None
    ) -> ExtractedContent:
        r"""
        在沙箱子进程中抽取附件内容；绝不抛出，失败一律降级为带 `note` 的 `unknown` 结果。

        Args:
            data: 附件的原始字节。
            file_metadata: 飞书附件元信息（`name` / `mime_type` 等），用于推断媒体类型与文件名。
            media_type: 媒体类型；省略时按魔数 / 元信息推断。默认为 `None`。

        Returns:
            [feishu.attachments.extractor.ExtractedContent][]。任何失败（超限 / 超时 / 解析异常）都返回
            `kind="unknown"` 并在 `note` 中给出中性英文说明，而不会抛出异常。
        """
        resolved = media_type or detect_media_type(data) or media_type_from_metadata(dict(file_metadata))
        if len(data) > self.max_bytes:
            return ExtractedContent(
                kind="unknown",
                media_type=resolved,
                size_bytes=len(data),
                note=f"attachment too large ({len(data)} bytes > {self.max_bytes} limit)",
            )
        name = _attachment_name(dict(file_metadata))
        async with self._semaphore:
            try:
                return await asyncio.wait_for(
                    asyncio.to_thread(_run_killable_extract, data, resolved, name, self.limits, self.timeout_seconds),
                    timeout=self.timeout_seconds + 10,  # backstop strictly exceeds the inner kill-and-cleanup budget
                )
            except asyncio.TimeoutError:
                return ExtractedContent(
                    kind="unknown", media_type=resolved, size_bytes=len(data), note="extraction timed out"
                )
            except Exception as exc:  # noqa: BLE001 — extraction must never raise into the caller
                return ExtractedContent(
                    kind="unknown",
                    media_type=resolved,
                    size_bytes=len(data),
                    note=f"extraction failed: {type(exc).__name__}",
                )


# --------------------------------------------------------------------------- #
# Sandbox plumbing (module-level so spawn can pickle the worker target).
# --------------------------------------------------------------------------- #


_MP_CONTEXT: Any = None


def _mp_context() -> Any:
    # Always spawn: this runs inside an asyncio worker thread, and forking a multithreaded process is
    # unsafe (deadlock-prone child on CPython 3.12+ — any lock held by another thread at fork stays
    # locked in the child). spawn starts a clean interpreter with no inherited locks.
    global _MP_CONTEXT
    if _MP_CONTEXT is None:
        _MP_CONTEXT = multiprocessing.get_context("spawn")
    return _MP_CONTEXT


def _run_killable_extract(
    data: bytes, media_type: str, name: str | None, limits: ExtractLimits, timeout_seconds: float
) -> ExtractedContent:
    ctx = _mp_context()
    result_queue: multiprocessing.Queue[Any] = ctx.Queue(maxsize=1)
    process = ctx.Process(target=_extract_worker, args=(result_queue, data, media_type, name, limits), daemon=True)
    process.start()
    try:
        status, payload = result_queue.get(timeout=max(1.0, timeout_seconds))
    except queue.Empty as exc:
        _stop_process(process)
        raise TimeoutError("attachment extraction timed out") from exc
    finally:
        process.join(timeout=1)
        if process.is_alive():
            _stop_process(process)
        result_queue.cancel_join_thread()  # never block the worker thread on a killed child's feeder
        result_queue.close()
    if status == "ok":
        return payload
    raise RuntimeError(str(payload))


def _stop_process(process: multiprocessing.Process) -> None:
    process.terminate()
    process.join(timeout=1)
    if process.is_alive():
        process.kill()
        process.join(timeout=1)


def _extract_worker(
    result_queue: multiprocessing.Queue[Any], data: bytes, media_type: str, name: str | None, limits: ExtractLimits
) -> None:
    try:
        result = extract_content(data, media_type=media_type, name=name, limits=limits)
    except Exception as exc:  # noqa: BLE001 — report any failure back to the parent (subprocess worker)
        result_queue.put(("error", f"{type(exc).__name__}: {exc}"))
    else:
        result_queue.put(("ok", result))


# --------------------------------------------------------------------------- #
# Pure (sync) extraction — runs inside the sandbox worker.
# --------------------------------------------------------------------------- #


def extract_content(data: bytes, *, media_type: str, name: str | None, limits: ExtractLimits) -> ExtractedContent:
    r"""同步抽取附件内容（在沙箱子进程内运行）；按媒体类型分派到图片 / PDF / 文本 / Office 提取。"""
    size = len(data)
    if media_type in _IMAGE_MIME_TYPES:
        if len(data) > limits.image_max_bytes:
            return ExtractedContent(
                kind="image", media_type=media_type, size_bytes=size, note="image exceeds the byte limit; omitted"
            )
        dims = image_dimensions(data, media_type) or _pil_dimensions(data, limits)
        if dims is None:  # fail closed: never forward an image whose dimensions we could not verify
            return ExtractedContent(
                kind="image",
                media_type=media_type,
                size_bytes=size,
                note="image dimensions could not be verified; omitted",
            )
        width, height = dims
        if width <= 0 or height <= 0 or width * height > limits.image_max_pixels:
            return ExtractedContent(
                kind="image", media_type=media_type, size_bytes=size, note="image exceeds the pixel limit; omitted"
            )
        return ExtractedContent(
            kind="image",
            media_type=media_type,
            size_bytes=size,
            images=[ExtractedImage(data, media_type, width, height)],
        )

    if _is_pdf(media_type, name):
        text, images, total_pages, truncated = extract_pdf_content(data, limits)
        return ExtractedContent(
            kind="pdf",
            media_type=media_type,
            size_bytes=size,
            text=text,
            total_pages=total_pages,
            truncated=truncated,
            images=[ExtractedImage(image, "image/png") for image in images],
            note=None if (text or images) else "no text or page images could be extracted",
        )

    text, kind, truncated = extract_text_like(data, media_type=media_type, name=name, limits=limits)
    if text is not None:
        return ExtractedContent(kind=kind, media_type=media_type, size_bytes=size, text=text, truncated=truncated)
    return ExtractedContent(
        kind="unknown",
        media_type=media_type,
        size_bytes=size,
        note="unsupported binary format with no extractable text",
    )


def extract_text_like(
    data: bytes, *, media_type: str, name: str | None, limits: ExtractLimits
) -> tuple[str | None, str, bool]:
    decoded = _decode_text(data, media_type=media_type, name=name)
    if decoded is not None:
        truncated_text, truncated = _truncate(decoded, limits)
        return truncated_text, "text", truncated
    suffix = _suffix(name)
    if suffix == ".docx" and _valid_office_zip(data, limits):
        return _finish_office(extract_docx_text(data, limits), "document", limits)
    if suffix == ".xlsx" and _valid_office_zip(data, limits):
        return _finish_office(extract_xlsx_text(data, limits), "spreadsheet", limits)
    if suffix == ".pptx" and _valid_office_zip(data, limits):
        return _finish_office(extract_pptx_text(data, limits), "presentation", limits)
    return None, "unknown", False


def _finish_office(text: str | None, kind: str, limits: ExtractLimits) -> tuple[str | None, str, bool]:
    if text is None:
        return None, kind, False
    truncated_text, truncated = _truncate(text, limits)
    return truncated_text, kind, truncated


def extract_pdf_content(data: bytes, limits: ExtractLimits) -> tuple[str | None, list[bytes], int, bool]:
    try:
        import fitz
    except ImportError:
        return None, [], 0, False
    try:
        document = fitz.open(stream=data, filetype="pdf")
    except Exception:
        return None, [], 0, False
    text_parts: list[str] = []
    images: list[bytes] = []
    truncated = False
    total_pages = len(document)
    try:
        current = 0
        for page_index in range(min(len(document), limits.pdf_extract_max_pages)):
            page_text = document[page_index].get_text("text").strip()
            if not page_text:
                continue
            remaining = limits.extract_text_max_chars - current
            if remaining <= 0:
                truncated = True
                break
            clipped = page_text[:remaining]
            text_parts.append(clipped)
            current += len(clipped)
            if len(page_text) > remaining:  # this page was cut off at the budget
                truncated = True
                break
        matrix = fitz.Matrix(limits.pdf_render_zoom, limits.pdf_render_zoom)
        for page_index in range(min(len(document), limits.pdf_render_max_pages)):
            page = document[page_index]
            width = int(page.rect.width * limits.pdf_render_zoom)
            height = int(page.rect.height * limits.pdf_render_zoom)
            if width <= 0 or height <= 0 or width * height > limits.pdf_render_max_pixels:
                continue
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            if pixmap.width * pixmap.height > limits.pdf_render_max_pixels:
                continue
            image = pixmap.tobytes("png")
            if len(image) <= limits.pdf_render_max_image_bytes:
                images.append(image)
    finally:
        document.close()
    text = "\n\n".join(text_parts).strip()
    return text or None, images, total_pages, truncated


def extract_docx_text(data: bytes, limits: ExtractLimits) -> str | None:
    try:
        from docx import Document
    except ImportError:
        return None
    try:
        document = Document(BytesIO(data))
    except Exception:
        return None
    lines: list[str] = []
    for paragraph in document.paragraphs:
        if _append_limited(lines, paragraph.text.strip(), limits):
            return "\n".join(lines).strip() or None
    for table in document.tables[: limits.docx_max_tables]:
        for row in table.rows[: limits.docx_max_table_rows]:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells and _append_limited(lines, " | ".join(cells), limits):
                return "\n".join(lines).strip() or None
    return "\n".join(lines).strip() or None


def extract_xlsx_text(data: bytes, limits: ExtractLimits) -> str | None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    try:
        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    except Exception:
        return None
    lines: list[str] = []
    try:
        for worksheet in workbook.worksheets[: limits.xlsx_max_sheets]:
            if _append_limited(lines, f"## {worksheet.title}", limits):
                break
            for row in worksheet.iter_rows(max_row=limits.xlsx_max_rows_per_sheet, values_only=True):
                values = [str(value) for value in row if value is not None and str(value).strip()]
                if values and _append_limited(lines, "\t".join(values), limits):
                    break
    finally:
        workbook.close()
    return "\n".join(lines).strip() or None


def extract_pptx_text(data: bytes, limits: ExtractLimits) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        presentation = Presentation(BytesIO(data))
    except Exception:
        return None
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        if index > limits.pptx_max_slides:
            break
        if _append_limited(lines, f"## Slide {index}", limits):
            break
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip() and _append_limited(lines, text.strip(), limits):
                break
    return "\n".join(lines).strip() or None


# --------------------------------------------------------------------------- #
# Detection, dimensions, guards, helpers.
# --------------------------------------------------------------------------- #


def detect_media_type(data: bytes) -> str | None:
    r"""按魔数（magic bytes）识别常见图片 / PDF 类型；无法识别返回 `None`。"""
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"GIF87a") or data.startswith(b"GIF89a"):
        return "image/gif"
    if data.startswith(b"%PDF"):
        return "application/pdf"
    if len(data) >= 12 and data[0:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    return None


def media_type_from_metadata(file_metadata: Mapping[str, Any]) -> str:
    r"""从飞书附件元信息（`mime_type` / 文件名）推断媒体类型；未知时返回 `application/octet-stream`。"""
    mime_type = file_metadata.get("mime_type")
    if isinstance(mime_type, str) and "/" in mime_type:
        return mime_type
    name = _attachment_name(file_metadata)
    if name:
        guessed, _ = mimetypes.guess_type(name)
        if guessed:
            return guessed
    return "application/octet-stream"


def image_dimensions(data: bytes, media_type: str) -> tuple[int, int] | None:
    r"""仅从文件头解析图片像素尺寸（不整图解码），用于像素炸弹防护。"""
    if media_type == "image/png" and len(data) >= 24 and data.startswith(b"\x89PNG\r\n\x1a\n"):
        return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")
    if media_type == "image/gif" and len(data) >= 10 and (data.startswith(b"GIF87a") or data.startswith(b"GIF89a")):
        return int.from_bytes(data[6:8], "little"), int.from_bytes(data[8:10], "little")
    if media_type == "image/jpeg" and data.startswith(b"\xff\xd8"):
        return _jpeg_dimensions(data)
    if media_type == "image/webp" and len(data) >= 30 and data[0:4] == b"RIFF" and data[8:12] == b"WEBP":
        return _webp_dimensions(data)
    return None


def _pil_dimensions(data: bytes, limits: ExtractLimits) -> tuple[int, int] | None:
    # Fallback when the header parser can't read dimensions. Pillow reports size without a full decode;
    # MAX_IMAGE_PIXELS makes a later decode raise on decompression bombs.
    try:
        from PIL import Image
    except ImportError:
        return None
    try:
        Image.MAX_IMAGE_PIXELS = limits.image_max_pixels
        with Image.open(BytesIO(data)) as image:
            return image.size
    except Exception:
        return None


def _jpeg_dimensions(data: bytes) -> tuple[int, int] | None:
    index = 2
    while index + 9 < len(data):
        if data[index] != 0xFF:
            index += 1
            continue
        marker = data[index + 1]
        index += 2
        if marker in {0xD8, 0xD9}:
            continue
        if index + 2 > len(data):
            return None
        length = int.from_bytes(data[index : index + 2], "big")
        if length < 2 or index + length > len(data):
            return None
        if marker in {0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF}:
            if length >= 7:
                height = int.from_bytes(data[index + 3 : index + 5], "big")
                width = int.from_bytes(data[index + 5 : index + 7], "big")
                return width, height
            return None
        index += length
    return None


def _webp_dimensions(data: bytes) -> tuple[int, int] | None:
    chunk = data[12:16]
    if chunk == b"VP8X" and len(data) >= 30:
        return 1 + int.from_bytes(data[24:27], "little"), 1 + int.from_bytes(data[27:30], "little")
    if chunk == b"VP8 " and len(data) >= 30:
        start = data.find(b"\x9d\x01\x2a", 20)
        if start >= 0 and start + 7 <= len(data):
            width = int.from_bytes(data[start + 3 : start + 5], "little") & 0x3FFF
            height = int.from_bytes(data[start + 5 : start + 7], "little") & 0x3FFF
            return width, height
    if chunk == b"VP8L" and len(data) >= 25:
        bits = int.from_bytes(data[21:25], "little")
        return (bits & 0x3FFF) + 1, ((bits >> 14) & 0x3FFF) + 1
    return None


def _valid_office_zip(data: bytes, limits: ExtractLimits) -> bool:
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            infos = archive.infolist()
            if len(infos) > limits.office_zip_max_entries:
                return False
            total = 0
            for info in infos:
                total += info.file_size
                if total > limits.office_zip_max_uncompressed_bytes:
                    return False
    except (zipfile.BadZipFile, OSError):
        return False
    return True


def _decode_text(data: bytes, *, media_type: str, name: str | None) -> str | None:
    if not _is_text(media_type, name):
        return None
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return None


def _is_text(media_type: str, name: str | None) -> bool:
    if media_type.startswith("text/") or media_type in _TEXT_MIME_TYPES:
        return True
    if not name:
        return False
    return name.lower().endswith((".csv", ".json", ".jsonl", ".log", ".md", ".txt", ".xml", ".yaml", ".yml"))


def _is_pdf(media_type: str, name: str | None) -> bool:
    return media_type == "application/pdf" or _suffix(name) == ".pdf"


def _suffix(name: str | None) -> str:
    if not name or "." not in name:
        return ""
    return f".{name.lower().rsplit('.', 1)[-1]}"


def _attachment_name(file_metadata: Mapping[str, Any]) -> str | None:
    name = file_metadata.get("name")
    return name if isinstance(name, str) and name.strip() else None


def _truncate(text: str, limits: ExtractLimits) -> tuple[str, bool]:
    if len(text) <= limits.inline_text_max_chars:
        return text, False
    return text[: limits.inline_text_max_chars] + "\n\n[attachment text truncated]", True


def _append_limited(lines: list[str], text: str, limits: ExtractLimits) -> bool:
    if not text:
        return False
    remaining = limits.extract_text_max_chars - sum(len(item) + 1 for item in lines)
    if remaining <= 0:
        return True
    lines.append(text[:remaining])
    return len(text) > remaining


# --------------------------------------------------------------------------- #
# Converter: ExtractedContent -> OpenAI-compatible chat content parts.
# --------------------------------------------------------------------------- #


def _data_url(data: bytes, media_type: str) -> str:
    return f"data:{media_type};base64,{base64.b64encode(data).decode('ascii')}"


def to_openai_content_parts(
    content: ExtractedContent, *, prompt: str | None = None, text_label: str = "Attachment text"
) -> list[dict[str, Any]]:
    r"""
    把 [feishu.attachments.extractor.ExtractedContent][] 转为 OpenAI Chat Completions 的多模态 `content` 列表。

    `prompt`（产品提供的分析提示词）作为首个文本块；随后是中性的元信息、抽取文本与图片（data URL）。SDK 不
    内置任何分析提示词，措辞由调用方决定。

    Args:
        content: 提取结果。
        prompt: 置于最前的分析提示词（产品文案）。
        text_label: 抽取文本块的标签。默认为 `"Attachment text"`。

    Returns:
        可作为 `messages[].content` 传给 OpenAI 兼容多模态接口的部件列表。
    """
    parts: list[dict[str, Any]] = []
    if prompt:
        parts.append({"type": "text", "text": prompt})
    header = f"kind={content.kind}; media_type={content.media_type}; size_bytes={content.size_bytes}"
    if content.total_pages is not None:
        header += f"; total_pages={content.total_pages}"
    if content.truncated:
        header += "; text_truncated=true"
    if content.note:
        header += f"; note={content.note}"
    parts.append({"type": "text", "text": header})
    if content.text:
        parts.append({"type": "text", "text": f"{text_label}:\n{content.text}"})
    for image in content.images:
        parts.append({"type": "image_url", "image_url": {"url": _data_url(image.data, image.media_type)}})
    return parts
